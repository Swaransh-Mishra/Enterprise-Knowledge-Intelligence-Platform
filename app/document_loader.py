"""
Multi-Format Document Loader.

Supports:
- PDF
- DOCX
- TXT
- Markdown
- CSV
- Excel
- PowerPoint
- HTML
- JSON
"""

from pathlib import Path
import json

import fitz
import pandas as pd
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation


class DocumentLoader:

    SUPPORTED_EXTENSIONS = {
        ".pdf",
        ".docx",
        ".txt",
        ".md",
        ".csv",
        ".xlsx",
        ".pptx",
        ".html",
        ".json",
    }

    def load_document(
        self,
        file_path: str
    ) -> str:

        path = Path(file_path)

        extension = path.suffix.lower()

        # --------------------------------
        # Check file exists
        # --------------------------------

        if not path.exists():

            raise FileNotFoundError(
                f"File not found: {file_path}"
            )

        # --------------------------------
        # Check supported extension
        # --------------------------------

        if extension not in self.SUPPORTED_EXTENSIONS:

            raise ValueError(
                f"Unsupported file type: {extension}"
            )

        # --------------------------------
        # Check file is not empty
        # --------------------------------

        if path.stat().st_size == 0:

            raise ValueError(
                "The uploaded file is empty."
            )

        loaders = {

            ".pdf": self._load_pdf,

            ".docx": self._load_docx,

            ".txt": self._load_txt,

            ".md": self._load_txt,

            ".csv": self._load_csv,

            ".xlsx": self._load_excel,

            ".pptx": self._load_pptx,

            ".html": self._load_html,

            ".json": self._load_json,
        }

        try:

            text = loaders[extension](
                str(path)
            )

        except Exception as error:

            raise ValueError(
                f"Could not read the uploaded "
                f"{extension} file. "
                f"The file may be invalid or corrupted."
            ) from error

        # --------------------------------
        # Validate extracted text
        # --------------------------------

        if not isinstance(text, str):

            raise ValueError(
                "The document loader did not "
                "return valid text."
            )

        text = text.strip()

        if not text:

            raise ValueError(
                "No readable text was found "
                "in the uploaded document."
            )

        return text

    # --------------------------------
    # PDF
    # --------------------------------

    def _load_pdf(
        self,
        file_path: str
    ):

        text = ""

        pdf = fitz.open(file_path)

        try:

            for page in pdf:

                text += page.get_text()

        finally:

            pdf.close()

        return text

    # --------------------------------
    # DOCX
    # --------------------------------

    def _load_docx(
        self,
        file_path: str
    ):

        doc = Document(file_path)

        return "\n".join(
            paragraph.text
            for paragraph in doc.paragraphs
        )

    # --------------------------------
    # TXT / Markdown
    # --------------------------------

    def _load_txt(
        self,
        file_path: str
    ):

        with open(
            file_path,
            "r",
            encoding="utf-8"
        ) as file:

            return file.read()

    # --------------------------------
    # CSV
    # --------------------------------

    def _load_csv(
        self,
        file_path: str
    ):

        df = pd.read_csv(file_path)

        return df.to_string(
            index=False
        )

    # --------------------------------
    # Excel
    # --------------------------------

    def _load_excel(
        self,
        file_path: str
    ):

        df = pd.read_excel(file_path)

        return df.to_string(
            index=False
        )

    # --------------------------------
    # PowerPoint
    # --------------------------------

    def _load_pptx(
        self,
        file_path: str
    ):

        presentation = Presentation(
            file_path
        )

        text = []

        for slide in presentation.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    if shape.text.strip():

                        text.append(
                            shape.text
                        )

        return "\n".join(text)

    # --------------------------------
    # HTML
    # --------------------------------

    def _load_html(
        self,
        file_path: str
    ):

        with open(
            file_path,
            "r",
            encoding="utf-8"
        ) as file:

            soup = BeautifulSoup(
                file.read(),
                "html.parser"
            )

        return soup.get_text(
            separator="\n"
        )

    # --------------------------------
    # JSON
    # --------------------------------

    def _load_json(
        self,
        file_path: str
    ):

        with open(
            file_path,
            "r",
            encoding="utf-8"
        ) as file:

            data = json.load(file)

        return json.dumps(
            data,
            indent=2,
            ensure_ascii=False
        )